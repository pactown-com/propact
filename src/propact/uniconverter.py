"""Universal document converter supporting 15+ formats."""

import os
import re
import smtplib
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Union, Any
from dataclasses import dataclass
from email import message_from_binary_file
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email import encoders
import logging

# Try to import optional dependencies
try:
    from markitdown import MarkItDown
    HAS_MARKITDOWN = True
except ImportError:
    HAS_MARKITDOWN = False

try:
    import weasyprint
    HAS_WEASYPRINT = True
except ImportError:
    HAS_WEASYPRINT = False

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    from pptx.util import Inches
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    import markdown
    from markdown.extensions import tables, fenced_code, toc, codehilite
    HAS_MARKDOWN = True
except ImportError:
    HAS_MARKDOWN = False

try:
    import pypandoc
    HAS_PYPANDOC = True
except ImportError:
    HAS_PYPANDOC = False


@dataclass
class ConversionResult:
    """Result of a universal conversion operation."""
    success: bool
    content: Optional[str] = None
    output_path: Optional[Path] = None
    format: Optional[str] = None
    errors: List[str] = None
    warnings: List[str] = None
    metadata: Dict[str, Any] = None
    
    def __post_init__(self):
        if self.errors is None:
            self.errors = []
        if self.warnings is None:
            self.warnings = []
        if self.metadata is None:
            self.metadata = {}


@dataclass
class EmailConfig:
    """Configuration for email sending."""
    host: str
    port: int
    username: str
    password: str
    use_tls: bool = True
    from_email: Optional[str] = None


class UniConverter:
    """Universal document converter supporting 15+ formats."""
    
    # Supported formats and their handlers
    TO_MARKDOWN = {
        "pdf": "markitdown",
        "docx": "markitdown",
        "doc": "markitdown",
        "pptx": "markitdown",
        "ppt": "markitdown",
        "xlsx": "markitdown",
        "xls": "markitdown",
        "html": "html",
        "htm": "html",
        "eml": "email",
        "msg": "email",
        "odt": "pandoc",
        "ods": "pandoc",
        "odp": "pandoc",
        "epub": "pandoc",
        "rtf": "pandoc",
    }
    
    FROM_MARKDOWN = {
        "pdf": "weasyprint",
        "html": "markdown",
        "htm": "markdown",
        "docx": "docx",
        "pptx": "pptx",
        "xlsx": "openpyxl",
        "eml": "email",
        "odt": "pandoc",
        "ods": "pandoc",
        "odp": "pandoc",
        "epub": "pandoc",
        "rtf": "pandoc",
    }
    
    def __init__(self):
        """Initialize the universal converter."""
        self.logger = logging.getLogger(__name__)
        
        # Initialize markitdown if available
        if HAS_MARKITDOWN:
            self.markitdown = MarkItDown()
        else:
            self.markitdown = None
            
        # Initialize markdown processor
        if HAS_MARKDOWN:
            self.md = markdown.Markdown(
                extensions=['tables', 'fenced_code', 'toc', 'codehilite'],
                extension_configs={
                    'codehilite': {
                        'css_class': 'highlight',
                        'use_pygments': True
                    }
                }
            )
        else:
            self.md = None
    
    def to_markdown(self, input_path: Union[str, Path], 
                   input_format: Optional[str] = None) -> ConversionResult:
        """Convert any supported format to Markdown.
        
        Args:
            input_path: Path to input file
            input_format: Format of input file (auto-detected if None)
            
        Returns:
            ConversionResult with markdown content
        """
        input_path = Path(input_path)
        
        if not input_path.exists():
            return ConversionResult(
                success=False,
                errors=[f"Input file not found: {input_path}"]
            )
        
        # Detect format if not provided
        if input_format is None:
            input_format = input_path.suffix.lower().lstrip('.')
        
        if input_format not in self.TO_MARKDOWN:
            return ConversionResult(
                success=False,
                errors=[f"Unsupported input format: {input_format}"]
            )
        
        handler = self.TO_MARKDOWN[input_format]
        
        try:
            if handler == "markitdown":
                return self._markitdown_to_md(input_path)
            elif handler == "html":
                return self._html_to_md(input_path)
            elif handler == "email":
                return self._email_to_md(input_path)
            elif handler == "pandoc":
                return self._pandoc_to_md(input_path, input_format)
            else:
                return ConversionResult(
                    success=False,
                    errors=[f"Unknown handler: {handler}"]
                )
                
        except Exception as e:
            self.logger.error(f"Conversion to markdown failed: {e}")
            return ConversionResult(
                success=False,
                errors=[str(e)]
            )
    
    def from_markdown(self, markdown_content: str, output_path: Union[str, Path],
                     output_format: Optional[str] = None) -> ConversionResult:
        """Convert Markdown to any supported format.
        
        Args:
            markdown_content: Markdown content to convert
            output_path: Path for output file
            output_format: Output format (auto-detected if None)
            
        Returns:
            ConversionResult with output path
        """
        output_path = Path(output_path)
        
        # Detect format if not provided
        if output_format is None:
            output_format = output_path.suffix.lower().lstrip('.')
        
        if output_format not in self.FROM_MARKDOWN:
            return ConversionResult(
                success=False,
                errors=[f"Unsupported output format: {output_format}"]
            )
        
        handler = self.FROM_MARKDOWN[output_format]
        
        try:
            if handler == "weasyprint":
                return self._md_to_pdf(markdown_content, output_path)
            elif handler == "markdown":
                return self._md_to_html(markdown_content, output_path)
            elif handler == "docx":
                return self._md_to_docx(markdown_content, output_path)
            elif handler == "pptx":
                return self._md_to_pptx(markdown_content, output_path)
            elif handler == "openpyxl":
                return self._md_to_xlsx(markdown_content, output_path)
            elif handler == "email":
                return self._md_to_email(markdown_content, output_path)
            elif handler == "pandoc":
                return self._pandoc_from_md(markdown_content, output_path, output_format)
            else:
                return ConversionResult(
                    success=False,
                    errors=[f"Unknown handler: {handler}"]
                )
                
        except Exception as e:
            self.logger.error(f"Conversion from markdown failed: {e}")
            return ConversionResult(
                success=False,
                errors=[str(e)]
            )
    
    def send_email(self, markdown_content: str, to_emails: Union[str, List[str]],
                   subject: str, config: EmailConfig,
                   attachments: Optional[List[Path]] = None) -> ConversionResult:
        """Send Markdown content as rich HTML email.
        
        Args:
            markdown_content: Markdown content to send
            to_emails: Recipient email(s)
            subject: Email subject
            config: SMTP configuration
            attachments: Optional file attachments
            
        Returns:
            ConversionResult
        """
        try:
            # Create message
            msg = MIMEMultipart('alternative')
            msg['Subject'] = subject
            msg['From'] = config.from_email or config.username
            
            if isinstance(to_emails, str):
                to_emails = [to_emails]
            msg['To'] = ', '.join(to_emails)
            
            # Add plain text version
            plain_part = MIMEText(markdown_content, 'plain')
            msg.attach(plain_part)
            
            # Add HTML version
            if self.md:
                html_content = self.md.convert(markdown_content)
                html_part = MIMEText(html_content, 'html')
                msg.attach(html_part)
            
            # Add attachments
            if attachments:
                for attachment in attachments:
                    if attachment.exists():
                        with open(attachment, 'rb') as f:
                            part = MIMEBase('application', 'octet-stream')
                            part.set_payload(f.read())
                            encoders.encode_base64(part)
                            part.add_header(
                                'Content-Disposition',
                                f'attachment; filename= {attachment.name}'
                            )
                            msg.attach(part)
            
            # Send email
            with smtplib.SMTP(config.host, config.port) as server:
                if config.use_tls:
                    server.starttls()
                server.login(config.username, config.password)
                server.send_message(msg)
            
            return ConversionResult(
                success=True,
                metadata={'recipients': to_emails, 'subject': subject}
            )
            
        except Exception as e:
            self.logger.error(f"Email sending failed: {e}")
            return ConversionResult(
                success=False,
                errors=[str(e)]
            )
    
    def _markitdown_to_md(self, input_path: Path) -> ConversionResult:
        """Convert using markitdown (MS Office/PDF)."""
        if not HAS_MARKITDOWN:
            return ConversionResult(
                success=False,
                errors=["markitdown is required for this conversion"]
            )
        
        if not self.markitdown:
            self.markitdown = MarkItDown()
        
        result = self.markitdown.convert(str(input_path))
        
        # markitdown returns a result object with text_content
        if hasattr(result, 'text_content'):
            content = result.text_content
        elif isinstance(result, dict):
            content = result.get('text_content', result.get('markdown', ''))
        else:
            content = str(result)
        
        return ConversionResult(
            success=True,
            content=content,
            format='markdown'
        )
    
    def _html_to_md(self, input_path: Path) -> ConversionResult:
        """Convert HTML to Markdown."""
        if not HAS_MARKDOWN:
            return ConversionResult(
                success=False,
                errors=["markdown is required for HTML conversion"]
            )
        
        with open(input_path, 'r', encoding='utf-8') as f:
            html_content = f.read()
        
        # Use html2text if available, otherwise basic conversion
        try:
            import html2text
            h = html2text.HTML2Text()
            h.ignore_links = False
            content = h.handle(html_content)
        except ImportError:
            # Basic HTML to text conversion
            import re
            content = re.sub(r'<[^>]+>', '', html_content)
        
        return ConversionResult(
            success=True,
            content=content,
            format='markdown'
        )
    
    def _email_to_md(self, input_path: Path) -> ConversionResult:
        """Convert email file to Markdown."""
        with open(input_path, 'rb') as f:
            msg = message_from_binary_file(f)
        
        content_parts = []
        
        # Extract headers
        content_parts.append(f"# Email Message")
        content_parts.append(f"**From:** {msg.get('From', 'Unknown')}")
        content_parts.append(f"**To:** {msg.get('To', 'Unknown')}")
        content_parts.append(f"**Subject:** {msg.get('Subject', 'No Subject')}")
        content_parts.append(f"**Date:** {msg.get('Date', 'Unknown')}")
        content_parts.append("")
        
        # Extract body
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() in ['text/plain', 'text/html']:
                    payload = part.get_payload(decode=True)
                    if payload:
                        charset = part.get_content_charset() or 'utf-8'
                        body = payload.decode(charset, errors='ignore')
                        
                        if part.get_content_type() == 'text/html' and HAS_MARKDOWN:
                            # Convert HTML to markdown
                            try:
                                import html2text
                                h = html2text.HTML2Text()
                                body = h.handle(body)
                            except ImportError:
                                pass
                        
                        content_parts.append(body)
        else:
            payload = msg.get_payload(decode=True)
            if payload:
                charset = msg.get_content_charset() or 'utf-8'
                body = payload.decode(charset, errors='ignore')
                content_parts.append(body)
        
        return ConversionResult(
            success=True,
            content='\n'.join(content_parts),
            format='markdown'
        )
    
    def _pandoc_to_md(self, input_path: Path, input_format: str) -> ConversionResult:
        """Convert using pandoc."""
        if not HAS_PYPANDOC:
            return ConversionResult(
                success=False,
                errors=["pypandoc is required for this conversion"]
            )
        
        try:
            content = pypandoc.convert_file(
                str(input_path),
                'markdown',
                format=input_format
            )
            
            return ConversionResult(
                success=True,
                content=content,
                format='markdown'
            )
        except Exception as e:
            return ConversionResult(
                success=False,
                errors=[f"Pandoc conversion failed: {str(e)}"]
            )
    
    def _md_to_pdf(self, markdown_content: str, output_path: Path) -> ConversionResult:
        """Convert Markdown to PDF using weasyprint."""
        if not HAS_WEASYPRINT:
            return ConversionResult(
                success=False,
                errors=["weasyprint is required for PDF generation"]
            )
        
        if not self.md:
            return ConversionResult(
                success=False,
                errors=["markdown is required for PDF generation"]
            )
        
        # Convert to HTML with CSS
        html_content = self.md.convert(markdown_content)
        
        # Add CSS styling
        css = """
        <style>
        body { font-family: Arial, sans-serif; margin: 2cm; }
        table { border-collapse: collapse; width: 100%; }
        th, td { border: 1px solid #ddd; padding: 8px; text-align: left; }
        th { background-color: #f2f2f2; }
        code { background-color: #f4f4f4; padding: 2px 4px; border-radius: 3px; }
        pre { background-color: #f4f4f4; padding: 10px; border-radius: 5px; overflow-x: auto; }
        </style>
        """
        
        full_html = f"<!DOCTYPE html><html><head>{css}</head><body>{html_content}</body></html>"
        
        # Generate PDF
        weasyprint.HTML(string=full_html).write_pdf(str(output_path))
        
        return ConversionResult(
            success=True,
            output_path=output_path,
            format='pdf'
        )
    
    def _md_to_html(self, markdown_content: str, output_path: Path) -> ConversionResult:
        """Convert Markdown to HTML."""
        if not self.md:
            return ConversionResult(
                success=False,
                errors=["markdown is required for HTML conversion"]
            )
        
        html_content = self.md.convert(markdown_content)
        
        # Add basic HTML structure
        full_html = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>Converted Document</title>
    <style>
        body {{ font-family: Arial, sans-serif; max-width: 800px; margin: 0 auto; padding: 20px; }}
        table {{ border-collapse: collapse; width: 100%; }}
        th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
        th {{ background-color: #f2f2f2; }}
        code {{ background-color: #f4f4f4; padding: 2px 4px; border-radius: 3px; }}
        pre {{ background-color: #f4f4f4; padding: 10px; border-radius: 5px; overflow-x: auto; }}
    </style>
</head>
<body>
{html_content}
</body>
</html>"""
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(full_html)
        
        return ConversionResult(
            success=True,
            output_path=output_path,
            format='html'
        )
    
    def _md_to_docx(self, markdown_content: str, output_path: Path) -> ConversionResult:
        """Convert Markdown to DOCX."""
        if not HAS_DOCX:
            return ConversionResult(
                success=False,
                errors=["python-docx is required for DOCX generation"]
            )
        
        doc = Document()
        
        # Parse markdown line by line
        lines = markdown_content.split('\n')
        current_table = []
        in_table = False
        
        for line in lines:
            if line.startswith('#'):
                # Header
                level = len(line) - len(line.lstrip('#'))
                text = line.lstrip('#').strip()
                doc.add_heading(text, level=level)
                in_table = False
            elif line.startswith('-') or line.startswith('*'):
                # List item
                p = doc.add_paragraph(line.strip()[1:].strip(), style='List Bullet')
                in_table = False
            elif '|' in line and line.count('|') >= 2:
                # Table row
                cells = [cell.strip() for cell in line.split('|') if cell.strip()]
                if len(cells) >= 2:
                    current_table.append(cells)
                    in_table = True
            elif not line.strip() and in_table and current_table:
                # End of table
                if len(current_table) >= 2:
                    # Add table to document
                    table = doc.add_table(rows=len(current_table), cols=len(current_table[0]))
                    table.style = 'Table Grid'
                    
                    for i, row in enumerate(current_table):
                        for j, cell in enumerate(row):
                            table.cell(i, j).text = cell
                    
                current_table = []
                in_table = False
            elif line.strip() and not in_table:
                # Regular paragraph
                doc.add_paragraph(line.strip())
        
        # Handle table at end of document
        if in_table and len(current_table) >= 2:
            table = doc.add_table(rows=len(current_table), cols=len(current_table[0]))
            table.style = 'Table Grid'
            
            for i, row in enumerate(current_table):
                for j, cell in enumerate(row):
                    table.cell(i, j).text = cell
        
        doc.save(str(output_path))
        
        return ConversionResult(
            success=True,
            output_path=output_path,
            format='docx'
        )
    
    def _md_to_pptx(self, markdown_content: str, output_path: Path) -> ConversionResult:
        """Convert Markdown to PPTX."""
        if not HAS_PPTX:
            return ConversionResult(
                success=False,
                errors=["python-pptx is required for PPTX generation"]
            )
        
        prs = Presentation()
        
        # Split content by headers for slides
        sections = re.split(r'\n(?=#+\s)', markdown_content)
        
        for section in sections:
            if not section.strip():
                continue
            
            lines = section.strip().split('\n')
            
            # Add slide
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title from first header
            if lines and lines[0].startswith('#'):
                title = lines[0].lstrip('#').strip()
                slide.shapes.title.text = title
                content_lines = lines[1:]
            else:
                slide.shapes.title.text = "Slide"
                content_lines = lines
            
            # Add content
            if content_lines:
                content_text = '\n'.join(content_lines)
                # Find content placeholder
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        shape.text_frame.text = content_text
                        break
        
        prs.save(str(output_path))
        
        return ConversionResult(
            success=True,
            output_path=output_path,
            format='pptx'
        )
    
    def _md_to_xlsx(self, markdown_content: str, output_path: Path) -> ConversionResult:
        """Convert Markdown tables to XLSX."""
        if not HAS_OPENPYXL:
            return ConversionResult(
                success=False,
                errors=["openpyxl is required for XLSX generation"]
            )
        
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Data"
        
        # Parse markdown for tables
        lines = markdown_content.split('\n')
        row_idx = 0
        
        for line in lines:
            if '|' in line and line.count('|') >= 2:
                # Table row
                cells = [cell.strip() for cell in line.split('|') if cell.strip()]
                
                # Skip separator lines
                if all(c.replace('-', '').replace(':', '').replace(' ', '') == '' for c in cells):
                    continue
                
                row_idx += 1
                for col_idx, cell in enumerate(cells, 1):
                    ws.cell(row=row_idx, column=col_idx, value=cell)
        
        wb.save(str(output_path))
        
        return ConversionResult(
            success=True,
            output_path=output_path,
            format='xlsx'
        )
    
    def _md_to_email(self, markdown_content: str, output_path: Path) -> ConversionResult:
        """Convert Markdown to email file (.eml)."""
        msg = MIMEMultipart('alternative')
        
        # Add headers
        msg['Subject'] = "Converted from Markdown"
        msg['From'] = "sender@example.com"
        msg['To'] = "recipient@example.com"
        
        # Add plain text version
        plain_part = MIMEText(markdown_content, 'plain')
        msg.attach(plain_part)
        
        # Add HTML version if markdown is available
        if self.md:
            html_content = self.md.convert(markdown_content)
            html_part = MIMEText(html_content, 'html')
            msg.attach(html_part)
        
        # Write to file
        with open(output_path, 'w') as f:
            f.write(msg.as_string())
        
        return ConversionResult(
            success=True,
            output_path=output_path,
            format='eml'
        )
    
    def _pandoc_from_md(self, markdown_content: str, output_path: Path,
                       output_format: str) -> ConversionResult:
        """Convert Markdown to other formats using pandoc."""
        if not HAS_PYPANDOC:
            return ConversionResult(
                success=False,
                errors=["pypandoc is required for this conversion"]
            )
        
        try:
            pypandoc.convert_text(
                markdown_content,
                output_format,
                format='markdown',
                outputfile=str(output_path)
            )
            
            return ConversionResult(
                success=True,
                output_path=output_path,
                format=output_format
            )
        except Exception as e:
            return ConversionResult(
                success=False,
                errors=[f"Pandoc conversion failed: {str(e)}"]
            )
    
    def get_supported_formats(self) -> Dict[str, List[str]]:
        """Get all supported formats.
        
        Returns:
            Dict with 'to_markdown' and 'from_markdown' keys
        """
        return {
            'to_markdown': list(self.TO_MARKDOWN.keys()),
            'from_markdown': list(self.FROM_MARKDOWN.keys())
        }


# Default converter instance
default_converter = UniConverter()
